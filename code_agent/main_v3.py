#!/usr/bin/env python3
"""
HTML to PowerPoint - Version 3: Maximum LLM Flexibility

Philosophy:
- Give LLM the HTML with MINIMAL constraints
- Trust LLM to analyze layout and create appropriate slide
- Provide powerful helper functions
- Let LLM decide positioning, styling, and structure

This approach scales better to new/unknown HTML layouts.
"""

import argparse
import json
import os
import re
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup
from databricks.sdk import WorkspaceClient
from pptx import Presentation

# Model configuration
CODE_GEN_ENDPOINT = "databricks-claude-sonnet-4-5"
model_serving_client = None


def load_html(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def capture_screenshot(html_path: str, output_path: str) -> bool:
    """Capture screenshot if requested"""
    try:
        from playwright.sync_api import sync_playwright
        
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            
            file_url = f"file://{os.path.abspath(html_path)}"
            page.goto(file_url)
            
            # Wait for charts
            try:
                page.wait_for_selector('canvas', timeout=5000)
                page.wait_for_timeout(1500)
            except:
                pass
            
            # Try to capture just the chart area
            canvas = page.query_selector('canvas')
            if canvas:
                chart_container = page.query_selector('canvas').evaluate_handle("""(canvas) => {
                    let elem = canvas;
                    for (let i = 0; i < 5; i++) {
                        if (elem.parentElement) {
                            elem = elem.parentElement;
                            const rect = elem.getBoundingClientRect();
                            if (elem.tagName === 'DIV' && rect.width > 300 && rect.height > 200) {
                                if (rect.width < window.innerWidth * 0.9) {
                                    return elem;
                                }
                            }
                        }
                    }
                    return canvas.parentElement || canvas;
                }""")
                
                if chart_container:
                    chart_container.as_element().screenshot(path=output_path)
                else:
                    canvas.screenshot(path=output_path)
            else:
                page.screenshot(path=output_path, full_page=False)
            
            browser.close()
            return True
            
    except Exception as e:
        print(f"⚠️  Screenshot failed: {e}")
        return False


# MINIMAL CONSTRAINT PROMPTS
SYSTEM_PROMPT = """You are a PowerPoint slide generator. Analyze HTML and create a professional slide.

Your task: Generate a Python module with convert_to_pptx(html_str, output_path, assets_dir) function.

Available tools (import from pptx):
- Presentation, slide_layouts (use layout 6 for blank)
- shapes.add_textbox(left, top, width, height) - positions in Inches()
- shapes.add_picture(path, left, top, width, height)
- shapes.add_chart(chart_type, left, top, width, height, chart_data)
- Pt() for font sizes, RGBColor() for colors, PP_ALIGN for alignment

Slide dimensions: 10" wide × 7.5" tall

Your approach:
1. Parse HTML to understand content and layout
2. Decide how to structure the slide (title, chart, stats, etc.)
3. Calculate appropriate positions for each element
4. Apply styling from HTML (colors, fonts, alignment)
5. Create a clean, professional layout

Key principles:
- Be adaptive - don't assume fixed layouts
- Extract visible content (titles, text, images, stats)
- Preserve colors and styling from HTML
- Use good visual hierarchy
- If screenshot.png exists in assets_dir, use it for charts
- If no screenshot but chart data exists in JavaScript, extract the data and create a PowerPoint chart using shapes.add_chart()
- For Chart.js data: extract labels and datasets, create appropriate chart type (bar, line, pie, etc.)
- Example Chart.js extraction: look for "const labels = Object.keys(lobTotals)" and "const data = Object.values(lobTotals)" patterns
- Use CategoryChartData to structure the data: chart_data.add_category(label) and chart_data.add_series('Series Name', data_values)
- Center-align titles if they appear centered in HTML
- Distribute elements evenly for clean appearance

Return ONLY the Python code, no explanations."""


USER_PROMPT_TEMPLATE = """Generate a PowerPoint slide from this HTML:

{html_content}

Screenshot mode: {use_screenshot}
{screenshot_note}

Instructions:
1. Analyze the HTML structure
2. Extract key content (title, subtitle, stats, charts)
3. If JavaScript contains chart data (like Chart.js), extract the data and create a PowerPoint chart
4. Design an appropriate layout
5. Generate Python code using python-pptx

For charts:
- If screenshot exists: use it
- If no screenshot but chart data in JavaScript: extract data and create PowerPoint chart
- Chart data format: look for arrays with labels and datasets, extract values and create appropriate chart type

Be creative and adaptive - create the best slide you can from this HTML.
Return only Python code with convert_to_pptx(html_str, output_path, assets_dir) function."""


def generate_slide_code(html_content: str, use_screenshot: bool, model: str) -> Optional[str]:
    """Ask LLM to generate slide code with maximum flexibility"""
    
    # Truncate HTML if too long
    if len(html_content) > 10000:
        soup = BeautifulSoup(html_content, 'lxml')
        # Keep body content, strip script tags
        for script in soup.find_all('script'):
            script.decompose()
        html_content = str(soup.body) if soup.body else html_content[:10000]
    
    screenshot_note = ""
    if use_screenshot:
        screenshot_note = """
If screenshot.png exists in assets_dir:
- Use it for the chart area: slide.shapes.add_picture(os.path.join(assets_dir, 'screenshot.png'), ...)
- Position it appropriately in your layout
- Don't recreate chart from data
"""
    
    prompt = USER_PROMPT_TEMPLATE.format(
        html_content=html_content,
        use_screenshot=use_screenshot,
        screenshot_note=screenshot_note
    )
    
    print(f"📤 Sending to LLM:")
    print(f"   • HTML length: {len(html_content)} chars")
    print(f"   • Prompt size: ~{len(prompt)} chars")
    
    try:
        response = model_serving_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,  # Some creativity but not too much
        )
        
        code = response.choices[0].message.content
        
        # Extract code from markdown if wrapped
        if '```python' in code:
            code = re.search(r'```python\n(.*?)```', code, re.DOTALL).group(1)
        elif '```' in code:
            code = re.search(r'```\n(.*?)```', code, re.DOTALL).group(1)
        
        print("✓ LLM generated code successfully")
        return code
        
    except Exception as e:
        print(f"❌ LLM call failed: {e}")
        return None


def main():
    global model_serving_client
    
    parser = argparse.ArgumentParser(description="HTML→PPT with Maximum LLM Flexibility")
    parser.add_argument("input_html", help="Path to HTML file")
    parser.add_argument("output_pptx", help="Path for output PowerPoint")
    parser.add_argument("--model", default=CODE_GEN_ENDPOINT, help="Model endpoint")
    parser.add_argument("--profile", default="logfood", help="Databricks profile")
    parser.add_argument("--use-screenshot", action="store_true", help="Capture chart as screenshot")
    args = parser.parse_args()
    
    print("="*80)
    print("HTML to PowerPoint - V3: Maximum LLM Flexibility")
    print("="*80)
    
    # Initialize Databricks
    print("\n🔧 Initializing Databricks client...")
    ws = WorkspaceClient(profile=args.profile, product='slide-generator')
    model_serving_client = ws.serving_endpoints.get_open_ai_client()
    print("✓ Connected")
    
    # Load HTML
    print(f"\n📄 Loading HTML: {args.input_html}")
    html_content = load_html(args.input_html)
    print("✓ Loaded")
    
    # Setup assets directory
    assets_dir = ".agent_build/assets_v3"
    os.makedirs(assets_dir, exist_ok=True)
    
    # Capture screenshot if requested
    if args.use_screenshot:
        print("\n📸 Capturing screenshot...")
        screenshot_path = os.path.join(assets_dir, "screenshot.png")
        if capture_screenshot(args.input_html, screenshot_path):
            print(f"✓ Screenshot saved: {screenshot_path}")
        else:
            print("⚠️  Screenshot failed, continuing without it")
            args.use_screenshot = False
    
    # Generate code with LLM
    print("\n" + "="*80)
    print("LLM CODE GENERATION (Flexible Approach)")
    print("="*80)
    
    code = generate_slide_code(html_content, args.use_screenshot, args.model)
    
    if not code:
        print("❌ Failed to generate code")
        return
    
    # Save generated module
    module_path = "generated_converter_v3.py"
    with open(module_path, 'w') as f:
        f.write(code)
    print(f"✓ Saved: {module_path}")
    
    # Execute generated code
    print("\n" + "="*80)
    print("EXECUTION")
    print("="*80)
    
    print("\n🔧 Loading generated module...")
    import importlib.util
    spec = importlib.util.spec_from_file_location("generated_converter_v3", module_path)
    module = importlib.util.module_from_spec(spec)
    
    try:
        spec.loader.exec_module(module)
        print("✓ Module loaded")
        
        print("\n🎨 Creating PowerPoint...")
        module.convert_to_pptx(html_content, args.output_pptx, assets_dir)
        
        # Verify file was created
        if os.path.exists(args.output_pptx):
            size = os.path.getsize(args.output_pptx) / 1024
            print(f"✓ Conversion complete")
            print(f"\n✅ SUCCESS: {args.output_pptx} ({size:.1f}KB)")
        else:
            print("❌ PowerPoint file was not created")
            
    except Exception as e:
        print(f"\n❌ Execution failed: {e}")
        import traceback
        traceback.print_exc()
    
    # Clean up screenshot if it was created
    if args.use_screenshot:
        screenshot_path = os.path.join(assets_dir, "screenshot.png")
        if os.path.exists(screenshot_path):
            try:
                os.remove(screenshot_path)
                print(f"\n🧹 Cleaned up screenshot: {screenshot_path}")
            except Exception as e:
                print(f"\n⚠️  Could not delete screenshot: {e}")
    
    print("\n" + "="*80)


if __name__ == "__main__":
    main()

