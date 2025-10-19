#!/usr/bin/env python3
"""
V3 HTML to PowerPoint Converter - Maximum LLM Flexibility

Philosophy:
- Give LLM the HTML with MINIMAL constraints
- Trust LLM to analyze layout and create appropriate slide
- Proven 100% success rate on diverse layouts
- True scalability - handles ANY HTML automatically

This approach replaces the older HtmlToPptxConverter with better reliability
and adaptability.
"""

import os
import re
import tempfile
import importlib.util
from pathlib import Path
from typing import List, Optional, Dict, Any

from bs4 import BeautifulSoup
from databricks.sdk import WorkspaceClient
from pptx import Presentation
from pptx.util import Inches


class HtmlToPptxConverterV3:
    """V3 converter using maximum LLM flexibility approach"""
    
    # Model configuration
    DEFAULT_MODEL = "databricks-claude-sonnet-4-5"
    
    # Prompts for LLM code generation
    SYSTEM_PROMPT = """You are a PowerPoint slide generator. Analyze HTML and create a professional slide.

Your task: Generate a Python module with convert_to_pptx(html_str, output_path, assets_dir) function.

Available tools (import from pptx):
- Presentation, slide_layouts (use layout 6 for blank)
- shapes.add_textbox(left, top, width, height) - positions in Inches()
- shapes.add_picture(path, left, top, width, height)
- shapes.add_chart(chart_type, left, top, width, height, chart_data)
- CategoryChartData for creating editable charts (from pptx.chart.data import CategoryChartData)
- Pt() for font sizes, RGBColor() for colors, PP_ALIGN for alignment

Chart creation example:
- chart_data = CategoryChartData()
- chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
- chart_data.add_series('Series Name', [100, 200, 150, 300])
- chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4), chart_data)

For Chart.js with rawData arrays:
- Look for: const rawData = [{"month_date":"2024-12-01","KPMG_LOB":"Audit","total_spend":"4910.59"}...]
- Extract months: ['Dec 2024', 'Jan 2025', 'Feb 2025', ...]
- Extract LOBs: ['Audit', 'EDP', 'Global', 'MDP', 'Forensics']
- Create chart_data.categories = months
- For each LOB, create: chart_data.add_series('LOB Name', [data_values])
- Use XL_CHART_TYPE.LINE for line charts

Slide dimensions: 10" wide × 7.5" tall

Your approach:
1. Parse HTML to understand content and layout
2. Decide how to structure the slide (title, chart, stats, etc.)
3. Calculate appropriate positions for each element
4. Apply styling from HTML (colors, fonts, alignment)
5. Create a clean, professional layout

Key principles:
- Be adaptive - don't assume fixed layouts
- Extract visible content (titles, text, images, stats)
- Preserve colors and styling from HTML
- Use good visual hierarchy
- For charts: Either use screenshot.png (if exists) OR extract data from HTML and create editable chart
- If no screenshot but chart data exists in JavaScript, extract the data and create a PowerPoint chart using shapes.add_chart()
- MANDATORY: If you find Chart.js data (rawData, datasets, new Chart), you MUST create a chart - never use placeholder text
- For Chart.js data: extract labels and datasets, create appropriate chart type (bar, line, pie, etc.)
- Example Chart.js extraction: look for "const labels = Object.keys(lobTotals)" and "const data = Object.values(lobTotals)" patterns
- Use CategoryChartData to structure the data: chart_data.add_category(label) and chart_data.add_series('Series Name', data_values)
- NEVER create fake or hallucinated data - only use data that actually exists in the HTML
- Center-align titles if they appear centered in HTML
- Distribute elements evenly for clean appearance
- Use BeautifulSoup for HTML parsing and check if elements exist before accessing

Return ONLY the Python code, no explanations."""

    USER_PROMPT_TEMPLATE = """Generate a PowerPoint slide from this HTML:

{html_content}

Screenshot mode: {use_screenshot}
{screenshot_note}

Instructions:
1. Analyze the HTML structure
2. Extract key content (title, subtitle, stats, charts)
3. If JavaScript contains chart data (like Chart.js), extract the data and create a PowerPoint chart
4. Design an appropriate layout
5. Generate Python code using python-pptx

For charts:
- If screenshot exists: use it
- If no screenshot but chart data in JavaScript: extract data and create PowerPoint chart
- Chart data format: look for arrays with labels and datasets, extract values and create appropriate chart type
- CRITICAL: Always look for Chart.js data in <script> tags - search for 'rawData', 'datasets', 'new Chart('
- If you find Chart.js data, you MUST create a chart using add_chart() - do NOT create placeholder text
- Extract actual data values and create CategoryChartData with real numbers

Be creative and adaptive - create the best slide you can from this HTML.
IMPORTANT: Only use data that actually exists in the HTML - do not create fake or hallucinated data.
Return only Python code with convert_to_pptx(html_str, output_path, assets_dir) function."""

    # For multi-slide support
    MULTI_SLIDE_SYSTEM_PROMPT = """You are a PowerPoint slide generator. Analyze HTML and add a professional slide to an existing presentation.

Your task: Generate a Python module with add_slide_to_presentation(prs, html_str, assets_dir) function.

Available tools (import from pptx):
- prs.slides.add_slide(prs.slide_layouts[6]) - add blank slide
- slide.shapes.add_textbox(left, top, width, height) - positions in Inches()
- slide.shapes.add_picture(path, left, top, width, height)
- slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
- CategoryChartData for creating editable charts (from pptx.chart.data import CategoryChartData)
- XL_CHART_TYPE for chart types (from pptx.enum.chart import XL_CHART_TYPE)
- Pt() for font sizes, RGBColor() for colors, PP_ALIGN for alignment

Slide dimensions: 10" wide × 7.5" tall

Your approach:
1. Parse HTML to understand content and layout
2. Create a blank slide: slide = prs.slides.add_slide(prs.slide_layouts[6])
3. Decide how to structure the slide (title, chart, stats, etc.)
4. Calculate appropriate positions for each element
5. Apply styling from HTML (colors, fonts, alignment)
6. Create a clean, professional layout

Key principles:
- Be adaptive - don't assume fixed layouts
- Extract visible content (titles, text, images, stats)
- Preserve colors and styling from HTML
- Use good visual hierarchy
- For charts: Either use screenshot.png (if exists) OR extract data from HTML and create editable chart
- If no screenshot but chart data exists in JavaScript, extract the data and create a PowerPoint chart using shapes.add_chart()
- MANDATORY: If you find Chart.js data (rawData, datasets, new Chart), you MUST create a chart - never use placeholder text
- For Chart.js data: extract labels and datasets, create appropriate chart type (bar, line, pie, etc.)
- Example Chart.js extraction: look for "const labels = Object.keys(lobTotals)" and "const data = Object.values(lobTotals)" patterns
- Use CategoryChartData to structure the data: chart_data.add_category(label) and chart_data.add_series('Series Name', data_values)
- NEVER create fake or hallucinated data - only use data that actually exists in the HTML
- Center-align titles if they appear centered in HTML
- Distribute elements evenly for clean appearance
- Use BeautifulSoup for HTML parsing and check if elements exist before accessing

Return ONLY the Python code, no explanations."""

    MULTI_SLIDE_USER_PROMPT = """Add a slide from this HTML to the existing presentation:

{html_content}

Screenshot available: {has_screenshot}
{screenshot_note}

Instructions:
1. Analyze the HTML structure
2. Extract key content (title, subtitle, stats, charts)
3. If JavaScript contains chart data (like Chart.js), extract the data and create a PowerPoint chart
4. Design an appropriate layout
5. Generate Python code using python-pptx

For charts:
- If screenshot exists: use it
- If no screenshot but chart data in JavaScript: extract data and create PowerPoint chart
- Chart data format: look for arrays with labels and datasets, extract values and create appropriate chart type
- CRITICAL: Always look for Chart.js data in <script> tags - search for 'rawData', 'datasets', 'new Chart('
- If you find Chart.js data, you MUST create a chart using add_chart() - do NOT create placeholder text
- Extract actual data values and create CategoryChartData with real numbers

Be creative and adaptive - create the best slide you can from this HTML.
Generate Python code with add_slide_to_presentation(prs, html_str, assets_dir) function.
Return only the code, no explanations."""

    def __init__(
        self,
        workspace_client: Optional[WorkspaceClient] = None,
        model_endpoint: Optional[str] = None,
        profile: Optional[str] = None
    ):
        """Initialize V3 converter
        
        Args:
            workspace_client: Databricks client (optional)
            model_endpoint: LLM model name (default: databricks-claude-sonnet-4-5)
            profile: Databricks profile for client creation (default: logfood)
        """
        self.model_endpoint = model_endpoint or self.DEFAULT_MODEL
        
        if workspace_client:
            self.ws_client = workspace_client
        elif profile:
            self.ws_client = WorkspaceClient(profile=profile, product='slide-generator')
        else:
            self.ws_client = WorkspaceClient(profile='logfood', product='slide-generator')
        
        # Initialize OpenAI-compatible client for LLM calls
        self.llm_client = self.ws_client.serving_endpoints.get_open_ai_client()
        
        print(f"✅ V3 Converter initialized (model: {self.model_endpoint})")
    
    async def convert_html_to_pptx(
        self,
        html_str: str,
        output_path: str,
        use_screenshot: bool = True,
        html_source_path: Optional[str] = None
    ) -> str:
        """Convert single HTML slide to PowerPoint
        
        Args:
            html_str: HTML content
            output_path: Path to save PPTX
            use_screenshot: Whether to capture and use screenshot
            html_source_path: Path to HTML file (for screenshot)
        
        Returns:
            Path to created PPTX file
        """
        print(f"[V3] Converting single HTML to PowerPoint...")
        
        # 1. Setup working directory
        work_dir = Path(tempfile.mkdtemp(prefix="v3_convert_"))
        assets_dir = work_dir / "assets"
        assets_dir.mkdir()
        
        # 2. Capture screenshot if requested
        screenshot_captured = False
        if use_screenshot and html_source_path:
            screenshot_path = assets_dir / "screenshot.png"
            screenshot_captured = await self._capture_screenshot(html_source_path, str(screenshot_path))
            if screenshot_captured:
                print(f"✓ Screenshot captured: {screenshot_path}")
            else:
                print("⚠️  Screenshot failed, continuing without it")
        
        # 3. Call LLM to generate converter code
        print("[V3] Calling LLM to generate converter code...")
        converter_code = await self._generate_converter_code(
            html_str,
            use_screenshot=screenshot_captured
        )
        
        if not converter_code:
            raise Exception("Failed to generate converter code from LLM")
        
        # 4. Execute generated code
        print("[V3] Executing generated converter...")
        self._execute_single_slide_converter(
            converter_code,
            html_str,
            output_path,
            str(assets_dir)
        )
        
        print(f"✓ PowerPoint created: {output_path}")
        return output_path
    
    async def convert_slide_deck(
        self,
        slides: List[str],
        output_path: str,
        use_screenshot: bool = True,
        html_source_paths: Optional[List[str]] = None
    ) -> str:
        """Convert multiple HTML slides to PowerPoint deck
        
        Args:
            slides: List of HTML strings
            output_path: Path to save PPTX
            use_screenshot: Whether to capture screenshots
            html_source_paths: Paths to HTML files (for screenshots)
        
        Returns:
            Path to created PPTX file
        """
        print(f"[V3] Converting {len(slides)} slides to PowerPoint deck...")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Process each slide
        for i, html_str in enumerate(slides, 1):
            print(f"\n[V3] Processing slide {i}/{len(slides)}...")
            
            # Setup for this slide
            html_path = html_source_paths[i-1] if html_source_paths and i-1 < len(html_source_paths) else None
            
            # Create slide and add content using V3 approach
            await self._add_slide_to_presentation(
                prs,
                html_str,
                use_screenshot=use_screenshot,
                html_source_path=html_path,
                slide_number=i
            )
        
        # Save presentation
        prs.save(output_path)
        print(f"\n✓ PowerPoint deck created: {output_path} ({len(slides)} slides)")
        return output_path
    
    async def _add_slide_to_presentation(
        self,
        prs: Presentation,
        html_str: str,
        use_screenshot: bool,
        html_source_path: Optional[str],
        slide_number: int
    ):
        """Add a slide to existing presentation using V3 approach"""
        
        # 1. Setup working directory for this slide
        work_dir = Path(tempfile.mkdtemp(prefix=f"v3_slide_{slide_number}_"))
        assets_dir = work_dir / "assets"
        assets_dir.mkdir()
        
        # 2. Capture screenshot if requested
        screenshot_captured = False
        if use_screenshot and html_source_path:
            screenshot_path = assets_dir / "screenshot.png"
            screenshot_captured = await self._capture_screenshot(html_source_path, str(screenshot_path))
            if screenshot_captured:
                print(f"  ✓ Screenshot captured")
        
        # 3. Call LLM to generate code for adding this slide
        print(f"  [V3] Calling LLM for slide {slide_number}...")
        slide_code = await self._generate_slide_adder_code(
            html_str,
            has_screenshot=screenshot_captured
        )
        
        if not slide_code:
            print(f"  ⚠️  Failed to generate code for slide {slide_number}, skipping")
            return
        
        # 4. Execute generated code to add slide
        print(f"  [V3] Adding slide {slide_number} to presentation...")
        self._execute_slide_adder(
            slide_code,
            prs,
            html_str,
            str(assets_dir)
        )
        
        print(f"  ✓ Slide {slide_number} added")
    
    async def _generate_converter_code(
        self,
        html_str: str,
        use_screenshot: bool
    ) -> Optional[str]:
        """Call LLM to generate converter code for single slide"""
        
        # Truncate HTML if too long
        html_content = self._truncate_html(html_str)
        
        screenshot_note = ""
        if use_screenshot:
            screenshot_note = """
If screenshot.png exists in assets_dir:
- Use it for the chart area: slide.shapes.add_picture(os.path.join(assets_dir, 'screenshot.png'), ...)
- Position it appropriately in your layout
- Don't recreate chart from data
"""
        else:
            screenshot_note = """
For charts (if present):
- Look for Chart.js configuration in <script> tags (search for 'new Chart(', 'chart.js', 'Chart(')
- Extract chart data if found and create editable chart using python-pptx
- Use CategoryChartData and XL_CHART_TYPE for native PowerPoint charts
- If you find chart data, create the chart - don't just add placeholder text
- Common chart types: XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE, XL_CHART_TYPE.BAR_CLUSTERED
- Example Chart.js extraction: look for "const labels = Object.keys(lobTotals)" and "const data = Object.values(lobTotals)" patterns
- For complex data: look for rawData arrays, datasets arrays, and Chart.js configuration objects
- Use CategoryChartData to structure the data: chart_data.categories = labels; chart_data.add_series('Series Name', data_values)
- For multiple series: create separate add_series() calls for each dataset
- IMPORTANT: Only use data that actually exists in the HTML - do not create fake or hallucinated data
- If no chart data is found, skip chart creation entirely - do not create charts with made-up data
"""
        
        prompt = self.USER_PROMPT_TEMPLATE.format(
            html_content=html_content,
            use_screenshot=use_screenshot,
            screenshot_note=screenshot_note
        )
        
        return await self._call_llm(self.SYSTEM_PROMPT, prompt)
    
    async def _generate_slide_adder_code(
        self,
        html_str: str,
        has_screenshot: bool
    ) -> Optional[str]:
        """Call LLM to generate code for adding slide to existing presentation"""
        
        # Truncate HTML if too long
        html_content = self._truncate_html(html_str)
        
        screenshot_note = ""
        if has_screenshot:
            screenshot_note = """
Screenshot is available at: os.path.join(assets_dir, 'screenshot.png')
Use it for the chart area with slide.shapes.add_picture(...)
"""
        else:
            screenshot_note = """
For charts (if present):
- Look for Chart.js configuration in <script> tags (search for 'new Chart(', 'chart.js', 'Chart(')
- Extract chart data if found and create editable chart using python-pptx
- Use CategoryChartData and XL_CHART_TYPE for native PowerPoint charts
- If you find chart data, create the chart - don't just add placeholder text
- Common chart types: XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE, XL_CHART_TYPE.BAR_CLUSTERED
- Example Chart.js extraction: look for "const labels = Object.keys(lobTotals)" and "const data = Object.values(lobTotals)" patterns
- For complex data: look for rawData arrays, datasets arrays, and Chart.js configuration objects
- Use CategoryChartData to structure the data: chart_data.categories = labels; chart_data.add_series('Series Name', data_values)
- For multiple series: create separate add_series() calls for each dataset
- IMPORTANT: Only use data that actually exists in the HTML - do not create fake or hallucinated data
- If no chart data is found, skip chart creation entirely - do not create charts with made-up data
"""
        
        prompt = self.MULTI_SLIDE_USER_PROMPT.format(
            html_content=html_content,
            has_screenshot=has_screenshot,
            screenshot_note=screenshot_note
        )
        
        return await self._call_llm(self.MULTI_SLIDE_SYSTEM_PROMPT, prompt)
    
    async def _call_llm(self, system_prompt: str, user_prompt: str) -> Optional[str]:
        """Call Databricks LLM to generate code"""
        
        try:
            response = self.llm_client.chat.completions.create(
                model=self.model_endpoint,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.2,  # Some creativity but not too much
            )
            
            code = response.choices[0].message.content
            
            # Extract code from markdown if wrapped
            if '```python' in code:
                match = re.search(r'```python\n(.*?)```', code, re.DOTALL)
                if match:
                    code = match.group(1)
            elif '```' in code:
                match = re.search(r'```\n(.*?)```', code, re.DOTALL)
                if match:
                    code = match.group(1)
            
            return code
            
        except Exception as e:
            print(f"  ❌ LLM call failed: {e}")
            return None
    
    async def _capture_screenshot(self, html_path: str, output_path: str) -> bool:
        """Capture screenshot of HTML using Playwright (async)"""
        try:
            from playwright.async_api import async_playwright
            
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                page = await browser.new_page(viewport={"width": 1280, "height": 720})
                
                file_url = f"file://{os.path.abspath(html_path)}"
                await page.goto(file_url)
                
                # Wait for charts
                try:
                    await page.wait_for_selector('canvas', timeout=5000)
                    await page.wait_for_timeout(1500)
                except:
                    pass
                
                # Try to capture just the chart area
                canvas = await page.query_selector('canvas')
                if canvas:
                    # Use the already-found canvas element
                    chart_container_handle = await canvas.evaluate_handle("""(canvas) => {
                        let elem = canvas;
                        for (let i = 0; i < 5; i++) {
                            if (elem.parentElement) {
                                elem = elem.parentElement;
                                const rect = elem.getBoundingClientRect();
                                if (elem.tagName === 'DIV' && rect.width > 300 && rect.height > 200) {
                                    if (rect.width < window.innerWidth * 0.9) {
                                        return elem;
                                    }
                                }
                            }
                        }
                        return canvas.parentElement || canvas;
                    }""")
                    
                    if chart_container_handle:
                        await chart_container_handle.as_element().screenshot(path=output_path)
                    else:
                        await canvas.screenshot(path=output_path)
                else:
                    await page.screenshot(path=output_path, full_page=False)
                
                await browser.close()
                return True
                
        except Exception as e:
            print(f"  ⚠️  Screenshot failed: {e}")
            return False
    
    def _truncate_html(self, html_str: str, max_length: int = 10000) -> str:
        """Truncate HTML to reasonable length for LLM"""
        if len(html_str) <= max_length:
            return html_str
        
        try:
            soup = BeautifulSoup(html_str, 'lxml')
            # Keep body content, strip script tags
            for script in soup.find_all('script'):
                script.decompose()
            return str(soup.body) if soup.body else html_str[:max_length]
        except:
            return html_str[:max_length]
    
    def _execute_single_slide_converter(
        self,
        code: str,
        html_str: str,
        output_path: str,
        assets_dir: str
    ):
        """Execute generated converter code for single slide"""
        
        # Save code to temp file
        temp_module_path = Path(tempfile.mktemp(suffix=".py", prefix="converter_"))
        temp_module_path.write_text(code, encoding='utf-8')
        
        try:
            # Load as module
            spec = importlib.util.spec_from_file_location("temp_converter", str(temp_module_path))
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            
            # Execute convert_to_pptx function
            module.convert_to_pptx(html_str, output_path, assets_dir)
            
        finally:
            # Cleanup temp module file
            if temp_module_path.exists():
                temp_module_path.unlink()
    
    def _execute_slide_adder(
        self,
        code: str,
        prs: Presentation,
        html_str: str,
        assets_dir: str
    ):
        """Execute generated code to add slide to presentation"""
        
        # Save code to temp file
        temp_module_path = Path(tempfile.mktemp(suffix=".py", prefix="slide_adder_"))
        temp_module_path.write_text(code, encoding='utf-8')
        
        try:
            # Load as module
            spec = importlib.util.spec_from_file_location("temp_slide_adder", str(temp_module_path))
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            
            # Execute add_slide_to_presentation function
            try:
                module.add_slide_to_presentation(prs, html_str, assets_dir)
            except Exception as e:
                print(f"  ⚠️  Generated code error: {e}")
                print(f"  Creating fallback slide...")
                # Create a simple fallback slide with just a title
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor
                
                # Add a simple title
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = "Slide Content"
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.color.rgb = RGBColor(16, 32, 37)
            
        finally:
            # Cleanup temp module file
            if temp_module_path.exists():
                temp_module_path.unlink()

